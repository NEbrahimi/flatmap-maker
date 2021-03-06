#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from collections import OrderedDict

#===============================================================================

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tqdm import tqdm

#===============================================================================

from mapmaker.properties import Properties

#===============================================================================

class Slide(object):
    def __init__(self, slide, properties, options):
        self.__properties = properties
        self.__names_only = options.names_only
        self.__slide_id = slide.slide_id
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            self.__notes = notes_slide.notes_text_frame.text.rstrip()
        else:
            self.__notes = ''
        self.__properties_by_id = OrderedDict()
        self.__external_ids = []
        self.__current_group = ['SLIDE']
        self.process_shape_list_(slide.shapes, True)

    def process_shape_list_(self, shapes, outermost):
        if outermost:
            progress_bar = tqdm(total=len(shapes),
                unit='shp', ncols=40,
                bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.__current_group.append(shape.name)
                self.process_shape_list_(shape.shapes, False)
                self.__current_group.pop()
            else:
                unique_id = '{}#{}'.format(self.__slide_id, shape.shape_id)
                self.__properties_by_id[unique_id] = self.__properties.get_properties(shape, self.__current_group[-1])
            if outermost:
                progress_bar.update(1)
        if outermost:
            progress_bar.close()

    def list(self):
        for id, properties in self.__properties_by_id.items():
            if properties and properties.get('shape-name', ' ')[0] in ['#', '.']:
                if self.__names_only:
                    print('{:8} {}'.format(id, properties['shape-name']))
                else:
                    print('{:8} {}'.format(id, properties))
        print('')

#===============================================================================

class Presentation(object):
    def __init__(self, powerpoint, properties, options):
        self.__pptx = pptx.Presentation(powerpoint)
        self.__slides_by_id = OrderedDict()
        self.__seen = []
        for slide in self.__pptx.slides:
            self.__slides_by_id[slide.slide_id] = Slide(slide, properties, options)

    def list(self):
        for slide in self.__slides_by_id.values():
            slide.list()

#===============================================================================

if __name__ == '__main__':
    import argparse
    import os, sys

    parser = argparse.ArgumentParser(description='Modify annotations in Powerpoint slides.')

#    parser.add_argument('-l', '--list', action='store_true',
#                        help='list shape annotations')

    parser.add_argument('--anatomical-map',
                        help='Excel spreadsheet file for mapping shape classes to anatomical entities')
    parser.add_argument('--properties',
                        help='JSON file specifying additional properties of shapes')

    parser.add_argument('-n', '--names-only', action='store_true',
                        help='Only list shape names')

    # Options `--replace`
    # specify slide
    #         shape by id/class

    #parser.add_argument('-r', '--replace', nargs=2, action='append',
    #                    help='find and replace text using ')

    parser.add_argument('powerpoint', help='Powerpoint file')

    args = parser.parse_args()
    args.label_database = 'labels.sqlite'
    external_properties = Properties(args)

    presentation = Presentation(args.powerpoint, external_properties, args)
    presentation.list()
